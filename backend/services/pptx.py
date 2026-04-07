from pptx import Presentation

def create_pptx(data: dict, output_path: str):

    prs = Presentation()

    for slide in data["slides"]:
        layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)

        if slide["number"] == 1:
            s.shapes.title.text = slide["title"]
            s.placeholders[1].text = slide["subtitle"]

        else:
            s.shapes.title.text = slide["name"]
            s.placeholders[1].text = slide["text"]

    prs.save(output_path)
